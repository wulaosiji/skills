#!/usr/bin/env python3
"""
BP (Business Plan) PPT Generator Template
This script generates a professional 10-page pitch deck using python-pptx.
Used as the reference implementation for the /business-plan-generator skill.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import sys

# ============================================================
# CONFIGURATION — Fill in project data here
# ============================================================

DATA = {
    "project_name": "示例项目",
    "tagline": "一句话定位：用AI重新定义行业效率",
    "core_values": [
        "10x 效率提升，替代传统人工流程",
        "基于大模型的智能决策引擎",
        "已服务100+企业客户"
    ],
    "date": "2026年4月",
    "team_label": "核心团队出品",

    # Page 2 - Market Pain Points
    "pain_points": [
        {"title": "效率低下", "desc": "行业仍依赖大量人工操作，平均耗时是AI方案的10倍"},
        {"title": "成本高企", "desc": "企业年均在此环节投入超50万元人力成本"},
        {"title": "质量不稳定", "desc": "人工处理错误率高达15%，严重影响业务决策"}
    ],
    "market_data": [
        {"label": "TAM", "value": "¥500亿", "desc": "全球市场规模"},
        {"label": "SAM", "value": "¥80亿", "desc": "中国可服务市场"},
        {"label": "SOM", "value": "¥5亿", "desc": "初期可获取市场"},
    ],

    # Page 3 - Solution
    "features": [
        {"title": "AI智能引擎", "desc": "基于大语言模型的自动化处理核心，准确率98%+"},
        {"title": "一键集成", "desc": "5分钟接入现有工作流，支持主流企业系统"},
        {"title": "实时监控", "desc": "全流程可视化，实时追踪处理状态和质量指标"},
        {"title": "数据安全", "desc": "企业级数据加密，支持私有化部署"}
    ],
    "differentiators": [
        "准确率行业领先（98% vs 行业平均85%）",
        "部署时间缩短90%（5分钟 vs 行业平均2周）",
        "成本降低70%以上"
    ],

    # Page 4 - Business Model
    "revenue_streams": [
        {"name": "SaaS订阅", "desc": "按月/年付费，分基础版/专业版/企业版", "pct": "60%"},
        {"name": "API调用", "desc": "按调用量计费，¥0.01/次", "pct": "25%"},
        {"name": "定制服务", "desc": "大客户定制化部署与咨询", "pct": "15%"}
    ],
    "pricing_tiers": [
        {"tier": "基础版", "price": "¥999/月", "target": "中小企业"},
        {"tier": "专业版", "price": "¥2,999/月", "target": "中型企业"},
        {"tier": "企业版", "price": "¥9,999/月", "target": "大型企业"}
    ],

    # Page 5 - Product Demo
    "product_status": "已上线Beta版，持续迭代中",
    "demo_features": [
        "智能仪表盘 — 一目了然的数据概览",
        "自动化工作流 — 拖拽式配置",
        "实时分析报告 — 自动生成洞察"
    ],

    # Page 6 - Competitive Analysis
    "competitors": [
        {"name": "竞品A", "f1": "✓", "f2": "✗", "f3": "✓", "f4": "✗", "f5": "✗"},
        {"name": "竞品B", "f1": "✓", "f2": "✓", "f3": "✗", "f4": "✗", "f5": "✗"},
        {"name": "竞品C", "f1": "✗", "f2": "✓", "f3": "✓", "f4": "✗", "f5": "✗"},
    ],
    "our_features": {"f1": "✓", "f2": "✓", "f3": "✓", "f4": "✓", "f5": "✓"},
    "feature_names": ["AI引擎", "一键集成", "实时监控", "私有部署", "API开放"],
    "moats": [
        "核心AI模型自研，准确率行业第一",
        "已积累100万条行业标注数据",
        "团队拥有10年行业经验"
    ],

    # Page 7 - Traction
    "milestones": [
        {"date": "2025.06", "event": "项目启动，核心团队组建"},
        {"date": "2025.09", "event": "MVP上线，首批10家种子用户"},
        {"date": "2025.12", "event": "产品迭代，用户增长至100+"},
        {"date": "2026.03", "event": "月营收突破¥50万，MoM增长40%"},
    ],
    "metrics": [
        {"label": "注册企业", "value": "120+"},
        {"label": "月活用户", "value": "3,500"},
        {"label": "月营收", "value": "¥50万"},
        {"label": "MoM增长", "value": "40%"}
    ],

    # Page 8 - Roadmap
    "roadmap": {
        "short": {"period": "0-6个月", "items": ["完善核心功能", "拓展至500家客户", "组建销售团队"]},
        "mid": {"period": "6-18个月", "items": ["推出企业版", "拓展东南亚市场", "年营收突破¥2000万"]},
        "long": {"period": "18个月+", "items": ["行业生态平台", "IPO准备", "全球化布局"]}
    },

    # Page 9 - Team
    "team": [
        {"name": "张三", "role": "CEO / 创始人", "bg": "前字节跳动技术总监，10年AI经验"},
        {"name": "李四", "role": "CTO", "bg": "前阿里云高级架构师，NLP方向博士"},
        {"name": "王五", "role": "COO", "bg": "前McKinsey顾问，5年企服销售经验"},
    ],
    "advisors": ["某知名VC合伙人", "某上市公司CEO"],

    # Page 10 - Fundraising
    "raise_amount": "¥2,000万",
    "equity": "10-15%",
    "valuation": "¥1.5亿（Pre-A轮）",
    "use_of_funds": [
        {"item": "研发投入", "pct": 50, "desc": "AI引擎升级 + 产品迭代"},
        {"item": "市场拓展", "pct": 30, "desc": "销售团队 + 品牌推广"},
        {"item": "运营 & 人才", "pct": 15, "desc": "核心人才招聘"},
        {"item": "储备资金", "pct": 5, "desc": "风险储备"}
    ],
    "contact": "contact@example.com"
}

# ============================================================
# DESIGN SYSTEM
# ============================================================

# Colors
PRIMARY = RGBColor(0x1A, 0x73, 0xE8)       # Tech Blue
DARK = RGBColor(0x20, 0x21, 0x24)           # Near Black
ACCENT = RGBColor(0x34, 0xA8, 0x53)         # Green
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)       # Light Gray
TEXT_DARK = RGBColor(0x20, 0x21, 0x24)
TEXT_LIGHT = RGBColor(0x5F, 0x63, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_ORANGE = RGBColor(0xFB, 0xBC, 0x04)  # Warm accent
ACCENT_RED = RGBColor(0xEA, 0x43, 0x35)     # Alert accent

# Fonts
FONT_CN = "Microsoft YaHei"
FONT_EN = "Calibri"

# Slide dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================
# HELPER FUNCTIONS
# ============================================================

def set_font(run, size=14, bold=False, color=TEXT_DARK, font_name=None):
    """Set font properties on a text run."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font_name:
        run.font.name = font_name


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name=None, anchor=MSO_ANCHOR.TOP):
    """Add a text box to a slide and return the shape."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = align
    tf.paragraphs[0].space_after = Pt(4)
    run = tf.paragraphs[0].add_run()
    run.text = text
    set_font(run, size, bold, color, font_name)
    try:
        txBox.text_frame.paragraphs[0].font.size = Pt(size)
    except:
        pass
    return txBox


def add_paragraph(text_frame, text, size=14, bold=False, color=TEXT_DARK,
                  align=PP_ALIGN.LEFT, space_before=0, space_after=6):
    """Add a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return p


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """Add a rectangle shape to the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_footer(slide, project_name, page_num, total=10):
    """Add footer bar with page number and project name."""
    # Bottom bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    bar.shadow.inherit = False

    # Project name on left
    add_textbox(slide, 0.5, 7.12, 4, 0.35, project_name, size=10, color=WHITE)

    # Page number on right
    add_textbox(slide, 11.5, 7.12, 1.5, 0.35, f"{page_num} / {total}",
                size=10, color=WHITE, align=PP_ALIGN.RIGHT)


def add_title_section(slide, title, subtitle=None, y_start=0.3):
    """Add a title area at the top of a slide."""
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y_start), Inches(0.08), Inches(0.6)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()
    line.shadow.inherit = False

    add_textbox(slide, 1.1, y_start - 0.05, 10, 0.7, title, size=30, bold=True, color=DARK)

    if subtitle:
        add_textbox(slide, 1.1, y_start + 0.6, 10, 0.4, subtitle, size=16, color=TEXT_LIGHT)

    return y_start + (1.2 if subtitle else 0.9)


def add_card(slide, left, top, width, height, title, body, accent_color=PRIMARY):
    """Add a styled card with colored top border."""
    # Card background
    card = add_rect(slide, left, top, width, height, fill_color=WHITE, border_color=RGBColor(0xE0, 0xE0, 0xE0), border_width=Pt(1))
    card.adjustments[0] = 0.02

    # Top accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    bar.shadow.inherit = False

    # Title
    add_textbox(slide, left + 0.2, top + 0.2, width - 0.4, 0.4, title, size=16, bold=True, color=DARK)

    # Body
    add_textbox(slide, left + 0.2, top + 0.65, width - 0.4, height - 0.9, body, size=13, color=TEXT_LIGHT)

    return card


def add_stat_block(slide, left, top, value, label, color=PRIMARY):
    """Add a large stat number with label below it."""
    add_textbox(slide, left, top, 2.5, 0.6, value, size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + 0.6, 2.5, 0.4, label, size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE GENERATORS
# ============================================================

def create_cover(prs, data):
    """Page 1: Cover / Project Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Full background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()
    bg.shadow.inherit = False

    # Decorative accent strip on left
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_H
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = PRIMARY
    strip.line.fill.background()
    strip.shadow.inherit = False

    # Project name
    add_textbox(slide, 1.5, 1.8, 10, 1.2, data["project_name"],
                size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Tagline
    add_textbox(slide, 1.5, 3.2, 10, 0.6, data["tagline"],
                size=22, color=RGBColor(0xA0, 0xC4, 0xF5), align=PP_ALIGN.LEFT)

    # Core values
    y = 4.2
    for val in data["core_values"]:
        # Bullet dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.7), Inches(y + 0.12), Inches(0.12), Inches(0.12)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        dot.shadow.inherit = False
        add_textbox(slide, 2.0, y, 9, 0.4, val, size=16, color=RGBColor(0xDA, 0xDA, 0xDA))
        y += 0.5

    # Bottom info
    add_textbox(slide, 1.5, 6.3, 5, 0.3, f"{data['team_label']}  |  {data['date']}",
                size=13, color=TEXT_LIGHT)

    # No standard footer on cover, add minimal page indicator
    add_textbox(slide, 11.5, 6.8, 1.5, 0.3, "1 / 10", size=10, color=TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def create_pain_points(prs, data):
    """Page 2: Market Pain Points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = add_title_section(slide, "市场痛点", "Market Pain Points — 真实刚需与行业现状")

    # Pain point cards
    colors = [ACCENT_RED, ACCENT_ORANGE, PRIMARY]
    card_width = 3.5
    gap = 0.4
    start_x = 0.8
    for i, pp in enumerate(data["pain_points"]):
        x = start_x + i * (card_width + gap)
        add_card(slide, x, y + 0.2, card_width, 2.2, pp["title"], pp["desc"], accent_color=colors[i % 3])

    # Market size stats
    stat_y = y + 2.9
    add_textbox(slide, 0.8, stat_y, 5, 0.4, "市场规模", size=18, bold=True, color=DARK)
    stat_y += 0.5

    stat_colors = [PRIMARY, ACCENT, ACCENT_ORANGE]
    for i, md in enumerate(data["market_data"]):
        x = 0.8 + i * 3.8
        # Stat background
        add_rect(slide, x, stat_y, 3.4, 1.5, fill_color=LIGHT_BG)
        add_stat_block(slide, x + 0.45, stat_y + 0.15, md["value"], md["label"], color=stat_colors[i])
        add_textbox(slide, x + 0.2, stat_y + 1.1, 3.0, 0.35, md["desc"], size=11, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    add_footer(slide, data["project_name"], 2)


def create_solution(prs, data):
    """Page 3: Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = add_title_section(slide, "解决方案", "Our Solution — 核心功能与差异化优势")

    # Feature cards (2x2 grid)
    card_w = 5.5
    card_h = 1.8
    gap_x = 0.5
    gap_y = 0.3
    for i, feat in enumerate(data["features"]):
        row = i // 2
        col = i % 2
        x = 0.8 + col * (card_w + gap_x)
        cy = y + 0.2 + row * (card_h + gap_y)
        add_card(slide, x, cy, card_w, card_h, feat["title"], feat["desc"], accent_color=PRIMARY)

    # Differentiators on the right side or bottom
    diff_y = y + 0.2 + 2 * (card_h + gap_y) + 0.1
    add_textbox(slide, 0.8, diff_y, 3, 0.35, "差异化优势", size=16, bold=True, color=ACCENT)
    for i, d in enumerate(data["differentiators"]):
        add_textbox(slide, 1.2, diff_y + 0.4 + i * 0.35, 10, 0.35,
                    f"▸ {d}", size=13, color=TEXT_DARK)

    add_footer(slide, data["project_name"], 3)


def create_business_model(prs, data):
    """Page 4: Business Model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = add_title_section(slide, "商业模式", "Business Model — 盈利方式与增长引擎")

    # Revenue streams
    stream_w = 3.5
    gap = 0.5
    for i, rs in enumerate(data["revenue_streams"]):
        x = 0.8 + i * (stream_w + gap)
        card = add_rect(slide, x, y + 0.2, stream_w, 2.4, fill_color=WHITE,
                        border_color=RGBColor(0xE0, 0xE0, 0xE0), border_width=Pt(1))
        # Percentage circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 1.2), Inches(y + 0.4), Inches(1.0), Inches(1.0)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY if i == 0 else (ACCENT if i == 1 else ACCENT_ORANGE)
        circle.line.fill.background()
        circle.shadow.inherit = False
        add_textbox(slide, x + 1.2, y + 0.65, 1.0, 0.5, rs["pct"],
                    size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, y + 1.55, stream_w - 0.4, 0.35, rs["name"],
                    size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, y + 1.95, stream_w - 0.4, 0.6, rs["desc"],
                    size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Pricing tiers table
    table_y = y + 3.0
    add_textbox(slide, 0.8, table_y, 3, 0.35, "定价策略", size=18, bold=True, color=DARK)
    table_y += 0.5

    cols = len(data["pricing_tiers"])
    tier_w = 3.5
    for i, tier in enumerate(data["pricing_tiers"]):
        x = 0.8 + i * (tier_w + 0.5)
        bg_color = PRIMARY if i == 1 else LIGHT_BG  # Highlight middle tier
        text_color = WHITE if i == 1 else DARK
        price_color = WHITE if i == 1 else PRIMARY

        rect = add_rect(slide, x, table_y, tier_w, 1.4, fill_color=bg_color)
        add_textbox(slide, x + 0.2, table_y + 0.1, tier_w - 0.4, 0.3, tier["tier"],
                    size=14, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, table_y + 0.45, tier_w - 0.4, 0.4, tier["price"],
                    size=24, bold=True, color=price_color, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, table_y + 0.95, tier_w - 0.4, 0.3, tier["target"],
                    size=12, color=text_color if i == 1 else TEXT_LIGHT, align=PP_ALIGN.CENTER)

    add_footer(slide, data["project_name"], 4)


def create_product_demo(prs, data):
    """Page 5: Product Demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = add_title_section(slide, "产品展示", "Product Demo — 产品形态与核心界面")

    # Status badge
    badge = add_rect(slide, 0.8, y + 0.1, 3.5, 0.4, fill_color=ACCENT)
    add_textbox(slide, 0.8, y + 0.12, 3.5, 0.35, f"  {data['product_status']}",
                size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Large mockup placeholder
    mock_x, mock_y = 0.8, y + 0.8
    mock_w, mock_h = 7.5, 4.5

    # Mockup frame
    frame = add_rect(slide, mock_x, mock_y, mock_w, mock_h,
                     fill_color=LIGHT_BG, border_color=RGBColor(0xD0, 0xD0, 0xD0), border_width=Pt(2))
    frame.adjustments[0] = 0.015

    # Browser-like top bar
    topbar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(mock_x), Inches(mock_y), Inches(mock_w), Inches(0.45)
    )
    topbar.fill.solid()
    topbar.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xED)
    topbar.line.fill.background()
    topbar.shadow.inherit = False

    # Browser dots
    for j, dc in enumerate([RGBColor(0xEA, 0x43, 0x35), RGBColor(0xFB, 0xBC, 0x04), RGBColor(0x34, 0xA8, 0x53)]):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(mock_x + 0.2 + j * 0.3), Inches(mock_y + 0.13), Inches(0.18), Inches(0.18)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = dc
        dot.line.fill.background()
        dot.shadow.inherit = False

    # Placeholder text
    add_textbox(slide, mock_x + 1.5, mock_y + 1.5, 4.5, 1.5,
                "[ 在此处插入产品截图 ]\nInsert Product Screenshot Here",
                size=20, color=RGBColor(0xB0, 0xB0, 0xB0), align=PP_ALIGN.CENTER)

    # Feature callouts on the right
    callout_x = mock_x + mock_w + 0.5
    callout_y = y + 0.8
    add_textbox(slide, callout_x, callout_y, 4, 0.4, "核心功能亮点", size=16, bold=True, color=DARK)
    callout_y += 0.5
    for i, feat in enumerate(data["demo_features"]):
        # Number circle
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(callout_x), Inches(callout_y + i * 1.2), Inches(0.4), Inches(0.4)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = PRIMARY
        num_circle.line.fill.background()
        num_circle.shadow.inherit = False
        add_textbox(slide, callout_x, callout_y + i * 1.2, 0.4, 0.4, str(i + 1),
                    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, callout_x + 0.55, callout_y + i * 1.2 - 0.05, 3.8, 0.8, feat,
                    size=13, color=TEXT_DARK)

    add_footer(slide, data["project_name"], 5)


def create_competitive_analysis(prs, data):
    """Page 6: Competitive Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = add_title_section(slide, "竞品分析", "Competitive Analysis — 竞争格局与核心壁垒")

    # Feature comparison table
    table_x, table_y = 0.8, y + 0.3
    n_features = len(data["feature_names"])
    col_w = 1.8
    row_h = 0.55
    name_col_w = 2.2

    # Header row
    add_rect(slide, table_x, table_y, name_col_w + n_features * col_w, row_h, fill_color=DARK)
    add_textbox(slide, table_x + 0.1, table_y + 0.05, name_col_w - 0.2, row_h - 0.1, "",
                size=13, bold=True, color=WHITE)
    for j, fn in enumerate(data["feature_names"]):
        add_textbox(slide, table_x + name_col_w + j * col_w, table_y + 0.05, col_w, row_h - 0.1,
                    fn, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Our row (highlighted)
    our_y = table_y + row_h
    add_rect(slide, table_x, our_y, name_col_w + n_features * col_w, row_h, fill_color=RGBColor(0xE8, 0xF0, 0xFE))
    add_textbox(slide, table_x + 0.1, our_y + 0.05, name_col_w - 0.2, row_h - 0.1,
                f"★ {data['project_name']}", size=13, bold=True, color=PRIMARY)
    for j in range(n_features):
        fkey = f"f{j+1}"
        val = data["our_features"].get(fkey, "✗")
        c = ACCENT if val == "✓" else ACCENT_RED
        add_textbox(slide, table_x + name_col_w + j * col_w, our_y + 0.05, col_w, row_h - 0.1,
                    val, size=16, bold=True, color=c, align=PP_ALIGN.CENTER)

    # Competitor rows
    for i, comp in enumerate(data["competitors"]):
        ry = our_y + (i + 1) * row_h
        bg = WHITE if i % 2 == 0 else LIGHT_BG
        add_rect(slide, table_x, ry, name_col_w + n_features * col_w, row_h, fill_color=bg)
        add_textbox(slide, table_x + 0.1, ry + 0.05, name_col_w - 0.2, row_h - 0.1,
                    comp["name"], size=13, color=TEXT_DARK)
        for j in range(n_features):
            fkey = f"f{j+1}"
            val = comp.get(fkey, "✗")
            c = ACCENT if val == "✓" else RGBColor(0xD0, 0xD0, 0xD0)
            add_textbox(slide, table_x + name_col_w + j * col_w, ry + 0.05, col_w, row_h - 0.1,
                        val, size=16, color=c, align=PP_ALIGN.CENTER)

    # Moats section
    moat_y = our_y + (len(data["competitors"]) + 1) * row_h + 0.5
    add_textbox(slide, 0.8, moat_y, 3, 0.4, "核心壁垒", size=18, bold=True, color=DARK)
    moat_y += 0.45
    for i, moat in enumerate(data["moats"]):
        # Shield icon placeholder
        shield = slide.shapes.add_shape(
            MSO_SHAPE.PENTAGON, Inches(1.0), Inches(moat_y + i * 0.6), Inches(0.3), Inches(0.35)
        )
        shield.fill.solid()
        shield.fill.fore_color.rgb = PRIMARY
        shield.line.fill.background()
        shield.shadow.inherit = False
        add_textbox(slide, 1.5, moat_y + i * 0.6 - 0.03, 10, 0.4, moat, size=14, color=TEXT_DARK)

    add_footer(slide, data["project_name"], 6)


def create_traction(prs, data):
    """Page 7: Traction / Progress"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = add_title_section(slide, "项目进展", "Traction — 关键里程碑与核心指标")

    # Key metrics row
    metric_y = y + 0.2
    metric_w = 2.7
    gap = 0.35
    for i, m in enumerate(data["metrics"]):
        x = 0.8 + i * (metric_w + gap)
        add_rect(slide, x, metric_y, metric_w, 1.4, fill_color=LIGHT_BG)
        colors = [PRIMARY, ACCENT, ACCENT_ORANGE, ACCENT_RED]
        add_stat_block(slide, x + 0.1, metric_y + 0.15, m["value"], m["label"], color=colors[i % 4])

    # Timeline
    timeline_y = metric_y + 1.9
    add_textbox(slide, 0.8, timeline_y, 3, 0.35, "发展历程", size=18, bold=True, color=DARK)
    timeline_y += 0.5

    # Horizontal timeline line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(timeline_y + 0.35), Inches(10), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()
    line.shadow.inherit = False

    n = len(data["milestones"])
    spacing = 10.0 / max(n - 1, 1)
    for i, ms in enumerate(data["milestones"]):
        mx = 1.5 + i * spacing
        # Dot on timeline
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(mx - 0.12), Inches(timeline_y + 0.22), Inches(0.3), Inches(0.3)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = PRIMARY
        dot.line.fill.background()
        dot.shadow.inherit = False
        # White inner dot
        inner = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(mx - 0.04), Inches(timeline_y + 0.3), Inches(0.14), Inches(0.14)
        )
        inner.fill.solid()
        inner.fill.fore_color.rgb = WHITE
        inner.line.fill.background()
        inner.shadow.inherit = False

        # Date above
        add_textbox(slide, mx - 0.8, timeline_y - 0.3, 1.8, 0.3, ms["date"],
                    size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        # Event below
        add_textbox(slide, mx - 0.8, timeline_y + 0.65, 1.8, 0.8, ms["event"],
                    size=11, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    add_footer(slide, data["project_name"], 7)


def create_roadmap(prs, data):
    """Page 8: Product Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = add_title_section(slide, "产品路线图", "Roadmap — 短期/中期/长期规划")

    phases = [
        ("short", "短期目标", PRIMARY),
        ("mid", "中期目标", ACCENT),
        ("long", "长期目标", ACCENT_ORANGE)
    ]

    phase_w = 3.5
    gap = 0.5
    for i, (key, label, color) in enumerate(phases):
        x = 0.8 + i * (phase_w + gap)
        phase_data = data["roadmap"][key]

        # Phase header
        header = add_rect(slide, x, y + 0.3, phase_w, 0.6, fill_color=color)
        add_textbox(slide, x, y + 0.35, phase_w, 0.5, label,
                    size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Period
        add_textbox(slide, x, y + 0.95, phase_w, 0.3, phase_data["period"],
                    size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

        # Items
        card_bg = add_rect(slide, x, y + 1.3, phase_w, 3.5, fill_color=LIGHT_BG)
        for j, item in enumerate(phase_data["items"]):
            # Checkmark
            add_textbox(slide, x + 0.15, y + 1.5 + j * 0.9, 0.3, 0.3, "→",
                        size=16, bold=True, color=color)
            add_textbox(slide, x + 0.5, y + 1.5 + j * 0.9, phase_w - 0.7, 0.8,
                        item, size=14, color=TEXT_DARK)

    # Arrow connectors between phases
    for i in range(2):
        ax = 0.8 + (i + 1) * (phase_w + gap) - gap / 2 - 0.15
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(y + 0.4), Inches(0.5), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
        arrow.line.fill.background()
        arrow.shadow.inherit = False

    add_footer(slide, data["project_name"], 8)


def create_team(prs, data):
    """Page 9: Team"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = add_title_section(slide, "团队介绍", "Team — 核心能力与资源背书")

    # Team member cards
    n = len(data["team"])
    card_w = 3.5
    gap = 0.5
    total_w = n * card_w + (n - 1) * gap
    start_x = (13.333 - total_w) / 2  # Center cards

    for i, member in enumerate(data["team"]):
        x = start_x + i * (card_w + gap)

        # Card background
        card = add_rect(slide, x, y + 0.3, card_w, 3.8, fill_color=WHITE,
                        border_color=RGBColor(0xE0, 0xE0, 0xE0), border_width=Pt(1))

        # Avatar placeholder (circle)
        avatar = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + card_w/2 - 0.5), Inches(y + 0.6), Inches(1.0), Inches(1.0)
        )
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = PRIMARY if i == 0 else (ACCENT if i == 1 else ACCENT_ORANGE)
        avatar.line.fill.background()
        avatar.shadow.inherit = False
        # Initials
        add_textbox(slide, x + card_w/2 - 0.5, y + 0.8, 1.0, 0.6, member["name"][0],
                    size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Name
        add_textbox(slide, x + 0.2, y + 1.75, card_w - 0.4, 0.4, member["name"],
                    size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

        # Role
        role_badge = add_rect(slide, x + 0.5, y + 2.2, card_w - 1.0, 0.35,
                              fill_color=RGBColor(0xE8, 0xF0, 0xFE))
        add_textbox(slide, x + 0.5, y + 2.22, card_w - 1.0, 0.3, member["role"],
                    size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

        # Background
        add_textbox(slide, x + 0.3, y + 2.75, card_w - 0.6, 1.2, member["bg"],
                    size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Advisors
    advisor_y = y + 4.4
    add_textbox(slide, 0.8, advisor_y, 3, 0.35, "顾问 & 资源", size=16, bold=True, color=DARK)
    advisor_text = "  |  ".join(data.get("advisors", []))
    add_textbox(slide, 0.8, advisor_y + 0.4, 11, 0.35, advisor_text, size=14, color=TEXT_LIGHT)

    add_footer(slide, data["project_name"], 9)


def create_fundraising(prs, data):
    """Page 10: Fundraising Plan"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = add_title_section(slide, "融资计划", "Fundraising — 融资金额与资金用途")

    # Key fundraising stats
    stats = [
        ("融资金额", data["raise_amount"]),
        ("出让比例", data["equity"]),
        ("估值", data["valuation"])
    ]
    stat_w = 3.5
    gap = 0.5
    for i, (label, value) in enumerate(stats):
        x = 0.8 + i * (stat_w + gap)
        bg_color = PRIMARY if i == 0 else LIGHT_BG
        text_color = WHITE if i == 0 else PRIMARY
        label_color = RGBColor(0xA0, 0xC4, 0xF5) if i == 0 else TEXT_LIGHT
        rect = add_rect(slide, x, y + 0.3, stat_w, 1.5, fill_color=bg_color)
        add_textbox(slide, x, y + 0.45, stat_w, 0.6, value,
                    size=32, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + 1.15, stat_w, 0.4, label,
                    size=14, color=label_color, align=PP_ALIGN.CENTER)

    # Use of Funds - Horizontal bar chart
    funds_y = y + 2.2
    add_textbox(slide, 0.8, funds_y, 3, 0.4, "资金用途", size=18, bold=True, color=DARK)
    funds_y += 0.5

    bar_colors = [PRIMARY, ACCENT, ACCENT_ORANGE, TEXT_LIGHT]
    max_bar_w = 7.0
    for i, fund in enumerate(data["use_of_funds"]):
        fy = funds_y + i * 0.85

        # Label
        add_textbox(slide, 0.8, fy, 2.5, 0.35, f"{fund['item']}  ({fund['pct']}%)",
                    size=13, bold=True, color=TEXT_DARK)

        # Bar background
        add_rect(slide, 3.5, fy + 0.05, max_bar_w, 0.4, fill_color=LIGHT_BG)

        # Bar fill
        bar_w = max_bar_w * fund["pct"] / 100
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(fy + 0.05), Inches(bar_w), Inches(0.4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_colors[i % 4]
        bar.line.fill.background()
        bar.shadow.inherit = False

        # Description
        add_textbox(slide, 3.5 + bar_w + 0.2, fy + 0.05, 4, 0.35, fund["desc"],
                    size=11, color=TEXT_LIGHT)

    # Contact
    contact_y = funds_y + len(data["use_of_funds"]) * 0.85 + 0.4
    contact_rect = add_rect(slide, 0.8, contact_y, 11.5, 0.6, fill_color=DARK)
    add_textbox(slide, 0.8, contact_y + 0.1, 11.5, 0.4,
                f"联系我们：{data['contact']}  |  期待与您合作",
                size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(slide, data["project_name"], 10)


# ============================================================
# MAIN
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    create_cover(prs, DATA)
    create_pain_points(prs, DATA)
    create_solution(prs, DATA)
    create_business_model(prs, DATA)
    create_product_demo(prs, DATA)
    create_competitive_analysis(prs, DATA)
    create_traction(prs, DATA)
    create_roadmap(prs, DATA)
    create_team(prs, DATA)
    create_fundraising(prs, DATA)

    output_path = os.path.join(os.getcwd(), f"{DATA['project_name']}_BP.pptx")
    prs.save(output_path)
    print(f"✅ BP PPT 已生成: {output_path}")
    print(f"共 {len(prs.slides)} 页")


if __name__ == "__main__":
    main()
