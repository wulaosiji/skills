#!/usr/bin/env python3
"""
BP to HTML Converter
Converts .pptx or .pdf Business Plan files into a self-contained HTML presentation.
Reference implementation for the /bp2html Claude Code skill.
"""

import os
import sys
import base64
import re
from html import escape

# ============================================================
# CONFIGURATION
# ============================================================
INPUT_FILE = ""  # Set this or pass as argument

# ============================================================
# PPTX EXTRACTION
# ============================================================

def extract_from_pptx(filepath):
    """Extract slide content from a .pptx file."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(filepath)
    slides = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "index": slide_idx,
            "texts": [],
            "images": [],
            "tables": [],
            "shapes": []
        }

        for shape in slide.shapes:
            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    img_bytes = image.blob
                    content_type = image.content_type
                    b64 = base64.b64encode(img_bytes).decode("utf-8")
                    slide_data["images"].append({
                        "data_uri": f"data:{content_type};base64,{b64}",
                        "width": shape.width,
                        "height": shape.height,
                        "left": shape.left,
                        "top": shape.top
                    })
                except Exception:
                    pass

            # Extract tables
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append(cell.text.strip())
                    rows.append(cells)
                slide_data["tables"].append(rows)

            # Extract text
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    # Determine text role by font size
                    font_size = None
                    is_bold = False
                    color_hex = None
                    for run in para.runs:
                        if run.font.size:
                            font_size = run.font.size.pt
                        if run.font.bold:
                            is_bold = True
                        if run.font.color and run.font.color.rgb:
                            color_hex = str(run.font.color.rgb)

                    role = "body"
                    if font_size:
                        if font_size >= 28:
                            role = "title"
                        elif font_size >= 18:
                            role = "subtitle"
                        elif font_size >= 15:
                            role = "heading"
                        elif font_size <= 11:
                            role = "caption"

                    slide_data["texts"].append({
                        "text": text,
                        "role": role,
                        "bold": is_bold,
                        "font_size": font_size,
                        "color": color_hex,
                        "left": shape.left,
                        "top": shape.top
                    })

            # Extract shape info (for colored blocks, etc.)
            if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
                shape_info = {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }
                try:
                    if shape.fill and shape.fill.fore_color and shape.fill.fore_color.rgb:
                        shape_info["fill_color"] = str(shape.fill.fore_color.rgb)
                except:
                    pass
                slide_data["shapes"].append(shape_info)

        # Sort texts by position (top to left)
        slide_data["texts"].sort(key=lambda t: (t.get("top", 0), t.get("left", 0)))
        slides.append(slide_data)

    return slides


# ============================================================
# PDF EXTRACTION
# ============================================================

def extract_from_pdf(filepath):
    """Extract page content from a .pdf file."""
    import fitz  # pymupdf

    doc = fitz.open(filepath)
    slides = []

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        slide_data = {
            "index": page_idx,
            "texts": [],
            "images": [],
            "tables": [],
            "shapes": []
        }

        # Extract text blocks
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if block["type"] == 0:  # Text block
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span["text"].strip()
                        if not text:
                            continue
                        font_size = span["size"]
                        is_bold = "bold" in span["font"].lower() or "heavy" in span["font"].lower()

                        role = "body"
                        if font_size >= 24:
                            role = "title"
                        elif font_size >= 16:
                            role = "subtitle"
                        elif font_size >= 13:
                            role = "heading"
                        elif font_size <= 9:
                            role = "caption"

                        color_int = span.get("color", 0)
                        color_hex = f"{color_int:06x}" if color_int else None

                        slide_data["texts"].append({
                            "text": text,
                            "role": role,
                            "bold": is_bold,
                            "font_size": font_size,
                            "color": color_hex,
                            "top": span["bbox"][1],
                            "left": span["bbox"][0]
                        })
            elif block["type"] == 1:  # Image block
                try:
                    img_data = block.get("image", None)
                    if img_data:
                        b64 = base64.b64encode(img_data).decode("utf-8")
                        ext = block.get("ext", "png")
                        mime = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                        slide_data["images"].append({
                            "data_uri": f"data:{mime};base64,{b64}",
                            "width": block["width"],
                            "height": block["height"]
                        })
                except:
                    pass

        # Extract images via get_images for better coverage
        for img_info in page.get_images(full=True):
            try:
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                if base_image:
                    b64 = base64.b64encode(base_image["image"]).decode("utf-8")
                    mime = base_image["ext"]
                    if mime == "jpg":
                        mime = "jpeg"
                    slide_data["images"].append({
                        "data_uri": f"data:image/{mime};base64,{b64}",
                        "width": base_image.get("width", 400),
                        "height": base_image.get("height", 300)
                    })
            except:
                pass

        slide_data["texts"].sort(key=lambda t: (t.get("top", 0), t.get("left", 0)))
        slides.append(slide_data)

    doc.close()
    return slides


# ============================================================
# SLIDE TYPE DETECTION
# ============================================================

BP_KEYWORDS = {
    "cover": ["bp", "商业计划", "business plan", "pitch deck"],
    "pain_points": ["痛点", "pain", "问题", "challenge", "市场", "market", "tam", "sam", "som"],
    "solution": ["解决方案", "solution", "功能", "feature", "产品", "product"],
    "business_model": ["商业模式", "business model", "盈利", "revenue", "定价", "pricing"],
    "demo": ["demo", "产品展示", "截图", "screenshot", "ui", "界面"],
    "competitive": ["竞品", "competitive", "竞争", "对比", "壁垒", "moat"],
    "traction": ["进展", "traction", "里程碑", "milestone", "指标", "metric"],
    "roadmap": ["路线图", "roadmap", "规划", "plan", "短期", "中期", "长期"],
    "team": ["团队", "team", "创始人", "founder", "ceo", "cto", "coo"],
    "fundraising": ["融资", "fundrais", "估值", "valuation", "资金", "fund"]
}

def detect_slide_type(slide_data, index, total):
    """Heuristically detect BP slide type."""
    if index == 0:
        return "cover"
    if index == total - 1:
        return "fundraising"

    all_text = " ".join([t["text"].lower() for t in slide_data["texts"]])

    best_type = "generic"
    best_score = 0
    for stype, keywords in BP_KEYWORDS.items():
        score = sum(1 for kw in keywords if kw in all_text)
        if score > best_score:
            best_score = score
            best_type = stype

    return best_type if best_score >= 1 else "generic"


# ============================================================
# HTML GENERATION
# ============================================================

SLIDE_TYPE_CLASSES = {
    "cover": "slide-cover",
    "pain_points": "slide-pain",
    "solution": "slide-solution",
    "business_model": "slide-business",
    "demo": "slide-demo",
    "competitive": "slide-competitive",
    "traction": "slide-traction",
    "roadmap": "slide-roadmap",
    "team": "slide-team",
    "fundraising": "slide-fundraising",
    "generic": "slide-generic"
}


def build_slide_html(slide_data, slide_type, index, total):
    """Build HTML for a single slide."""
    css_class = SLIDE_TYPE_CLASSES.get(slide_type, "slide-generic")
    parts = []
    parts.append(f'<div class="slide {css_class}" data-index="{index}">')
    parts.append('<div class="slide-inner">')

    # Group texts by role
    titles = [t for t in slide_data["texts"] if t["role"] == "title"]
    subtitles = [t for t in slide_data["texts"] if t["role"] == "subtitle"]
    headings = [t for t in slide_data["texts"] if t["role"] == "heading"]
    bodies = [t for t in slide_data["texts"] if t["role"] == "body"]
    captions = [t for t in slide_data["texts"] if t["role"] == "caption"]

    # Render titles
    for t in titles:
        text = escape(t["text"])
        style = ""
        if t.get("color") and t["color"] not in ("000000", "202124"):
            style = f' style="color:#{t["color"]}"'
        if slide_type == "cover":
            parts.append(f'<h1 class="slide-title cover-title"{style}>{text}</h1>')
        else:
            parts.append(f'<h2 class="slide-title"{style}>{text}</h2>')

    # Render subtitles
    for t in subtitles:
        text = escape(t["text"])
        parts.append(f'<p class="slide-subtitle">{text}</p>')

    # Render images
    if slide_data["images"]:
        parts.append('<div class="slide-images">')
        for img in slide_data["images"]:
            parts.append(f'<img src="{img["data_uri"]}" class="slide-img" alt="slide image" loading="lazy">')
        parts.append('</div>')

    # Render headings + body as content blocks
    if headings or bodies:
        parts.append('<div class="slide-content">')

        # Check if content looks like cards (multiple headings with body text)
        if len(headings) >= 2:
            parts.append('<div class="card-grid">')
            # Pair headings with following body texts
            heading_positions = [(h, h.get("top", 0)) for h in headings]
            body_positions = [(b, b.get("top", 0)) for b in bodies]

            for h, h_top in heading_positions:
                parts.append('<div class="card">')
                parts.append(f'<h3 class="card-title">{escape(h["text"])}</h3>')
                # Find body texts near this heading
                for b, b_top in body_positions:
                    if abs(b_top - h_top) < 200000:  # EMU proximity
                        parts.append(f'<p class="card-body">{escape(b["text"])}</p>')
                parts.append('</div>')
            parts.append('</div>')
        else:
            for h in headings:
                parts.append(f'<h3>{escape(h["text"])}</h3>')
            for b in bodies:
                text = escape(b["text"])
                if text.startswith(("▸", "•", "·", "-", "→", "✓", "✗", "★")):
                    parts.append(f'<li class="bullet">{text}</li>')
                elif b.get("bold"):
                    parts.append(f'<p><strong>{text}</strong></p>')
                else:
                    parts.append(f'<p>{text}</p>')

        parts.append('</div>')

    # Render tables
    for table_data in slide_data["tables"]:
        parts.append('<div class="table-wrapper">')
        parts.append('<table>')
        for row_idx, row in enumerate(table_data):
            tag = "th" if row_idx == 0 else "td"
            parts.append("<tr>")
            for cell in row:
                cell_text = escape(cell)
                css = ""
                if cell_text == "✓":
                    css = ' class="check"'
                elif cell_text == "✗":
                    css = ' class="cross"'
                parts.append(f"<{tag}{css}>{cell_text}</{tag}>")
            parts.append("</tr>")
        parts.append('</table>')
        parts.append('</div>')

    # Render captions
    for c in captions:
        parts.append(f'<p class="caption">{escape(c["text"])}</p>')

    # Slide number
    parts.append(f'<div class="slide-number">{index + 1}</div>')

    parts.append('</div>')  # slide-inner
    parts.append('</div>')  # slide

    return "\n".join(parts)


def generate_html(slides_data, title, accent_color="#1a73e8"):
    """Generate complete self-contained HTML presentation."""
    total = len(slides_data)

    # Build slide HTML
    slides_html = []
    for i, sd in enumerate(slides_data):
        stype = detect_slide_type(sd, i, total)
        slides_html.append(build_slide_html(sd, stype, i, total))

    slides_joined = "\n\n".join(slides_html)

    # Build dots
    dots = "\n".join(
        f'<span class="dot{" active" if i == 0 else ""}" data-index="{i}"></span>'
        for i in range(total)
    )

    html = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{escape(title)} — Business Plan</title>
<style>
*, *::before, *::after {{ margin:0; padding:0; box-sizing:border-box; }}
:root {{
  --accent: {accent_color};
  --accent-light: {accent_color}22;
  --dark: #202124;
  --text: #202124;
  --text2: #5f6368;
  --bg: #0f0f1a;
  --slide-bg: #ffffff;
  --radius: 12px;
}}
html {{ font-size: 16px; }}
body {{
  font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", "Microsoft YaHei", "PingFang SC", sans-serif;
  background: linear-gradient(135deg, #0f0f1a 0%, #1a1a2e 50%, #16213e 100%);
  color: var(--text);
  min-height: 100vh;
  overflow: hidden;
  user-select: none;
  -webkit-user-select: none;
}}

/* Progress Bar */
.progress-bar {{
  position: fixed; top:0; left:0; width:100%; height:3px;
  background: rgba(255,255,255,0.1); z-index:100;
}}
.progress-fill {{
  height:100%; width:0; background: var(--accent);
  transition: width 0.4s ease;
}}

/* Slides */
.slides-container {{
  position: relative;
  width: 100vw; height: 100vh;
  display: flex; align-items: center; justify-content: center;
}}
.slide {{
  position: absolute;
  width: min(90vw, 1200px);
  height: min(80vh, 675px);
  background: var(--slide-bg);
  border-radius: var(--radius);
  box-shadow: 0 20px 60px rgba(0,0,0,0.3), 0 0 0 1px rgba(255,255,255,0.05);
  opacity: 0;
  transform: translateX(80px) scale(0.95);
  transition: all 0.5s cubic-bezier(0.4, 0, 0.2, 1);
  pointer-events: none;
  overflow-y: auto;
  overflow-x: hidden;
}}
.slide.active {{
  opacity: 1;
  transform: translateX(0) scale(1);
  pointer-events: auto;
}}
.slide.prev {{
  opacity: 0;
  transform: translateX(-80px) scale(0.95);
}}
.slide-inner {{
  padding: 50px 60px;
  height: 100%;
  position: relative;
}}

/* Cover slide */
.slide-cover {{
  background: linear-gradient(135deg, var(--dark) 0%, #1a1a2e 100%);
  color: #ffffff;
}}
.slide-cover .slide-title {{ color: #ffffff; }}
.slide-cover .slide-subtitle {{ color: rgba(255,255,255,0.7); }}
.slide-cover .bullet {{ color: rgba(255,255,255,0.85); }}
.slide-cover .caption {{ color: rgba(255,255,255,0.5); }}
.cover-title {{
  font-size: 2.8rem;
  margin-bottom: 0.5em;
  line-height: 1.2;
}}

/* Titles */
.slide-title {{
  font-size: 1.8rem;
  font-weight: 700;
  color: var(--text);
  margin-bottom: 0.3em;
  padding-bottom: 0.3em;
  border-bottom: 3px solid var(--accent);
  display: inline-block;
}}
.slide-subtitle {{
  font-size: 1.05rem;
  color: var(--text2);
  margin-bottom: 1.5em;
}}

/* Content */
.slide-content {{
  margin-top: 0.8em;
  line-height: 1.7;
}}
.slide-content p {{
  margin-bottom: 0.5em;
  font-size: 0.95rem;
  color: var(--text);
}}
.slide-content h3 {{
  font-size: 1.1rem;
  color: var(--text);
  margin: 0.8em 0 0.3em;
}}
.bullet {{
  list-style: none;
  padding-left: 1.2em;
  position: relative;
  margin-bottom: 0.4em;
  font-size: 0.95rem;
}}
.bullet::before {{
  content: "";
  position: absolute;
  left: 0; top: 0.5em;
  width: 6px; height: 6px;
  background: var(--accent);
  border-radius: 50%;
}}

/* Card Grid */
.card-grid {{
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(240px, 1fr));
  gap: 16px;
  margin-top: 0.5em;
}}
.card {{
  background: #f8f9fa;
  border-radius: 8px;
  padding: 20px;
  border-top: 3px solid var(--accent);
  transition: transform 0.2s, box-shadow 0.2s;
}}
.card:hover {{
  transform: translateY(-2px);
  box-shadow: 0 8px 24px rgba(0,0,0,0.08);
}}
.card:nth-child(2) {{ border-top-color: #34a853; }}
.card:nth-child(3) {{ border-top-color: #fbbc04; }}
.card:nth-child(4) {{ border-top-color: #ea4335; }}
.card-title {{
  font-size: 1rem;
  font-weight: 700;
  margin-bottom: 0.4em;
  color: var(--text);
}}
.card-body {{
  font-size: 0.85rem;
  color: var(--text2);
  line-height: 1.5;
}}

/* Images */
.slide-images {{
  display: flex;
  gap: 16px;
  flex-wrap: wrap;
  justify-content: center;
  margin: 1em 0;
}}
.slide-img {{
  max-width: 100%;
  max-height: 350px;
  border-radius: 8px;
  object-fit: contain;
}}

/* Tables */
.table-wrapper {{
  overflow-x: auto;
  margin: 1em 0;
}}
table {{
  width: 100%;
  border-collapse: collapse;
  font-size: 0.85rem;
}}
th {{
  background: var(--dark);
  color: #fff;
  padding: 10px 14px;
  text-align: left;
  font-weight: 600;
}}
td {{
  padding: 8px 14px;
  border-bottom: 1px solid #e8eaed;
}}
tr:nth-child(even) td {{
  background: #f8f9fa;
}}
tr:hover td {{
  background: var(--accent-light);
}}
.check {{ color: #34a853; font-weight: bold; text-align: center; }}
.cross {{ color: #d0d0d0; text-align: center; }}

/* Caption */
.caption {{
  font-size: 0.75rem;
  color: var(--text2);
  margin-top: 0.3em;
}}

/* Slide Number */
.slide-number {{
  position: absolute;
  bottom: 16px; right: 24px;
  font-size: 0.7rem;
  color: var(--text2);
  opacity: 0.5;
}}
.slide-cover .slide-number {{ color: rgba(255,255,255,0.3); }}

/* Navigation */
.nav-dots {{
  position: fixed;
  bottom: 24px; left: 50%;
  transform: translateX(-50%);
  display: flex; gap: 8px;
  z-index: 50;
}}
.dot {{
  width: 10px; height: 10px;
  border-radius: 50%;
  background: rgba(255,255,255,0.3);
  cursor: pointer;
  transition: all 0.3s;
}}
.dot.active {{
  background: var(--accent);
  transform: scale(1.3);
}}
.dot:hover {{
  background: rgba(255,255,255,0.6);
}}
.slide-counter {{
  position: fixed;
  bottom: 24px; right: 24px;
  color: rgba(255,255,255,0.5);
  font-size: 0.85rem;
  z-index: 50;
}}
.nav-arrow {{
  position: fixed;
  top: 50%; transform: translateY(-50%);
  background: rgba(255,255,255,0.08);
  border: 1px solid rgba(255,255,255,0.1);
  color: rgba(255,255,255,0.6);
  width: 44px; height: 44px;
  border-radius: 50%;
  font-size: 1.4rem;
  cursor: pointer;
  z-index: 50;
  display: flex; align-items: center; justify-content: center;
  transition: all 0.2s;
  backdrop-filter: blur(4px);
}}
.nav-arrow:hover {{
  background: rgba(255,255,255,0.15);
  color: #fff;
}}
.nav-arrow.prev {{ left: 20px; }}
.nav-arrow.next {{ right: 20px; }}
.fullscreen-btn {{
  position: fixed;
  top: 20px; right: 20px;
  background: rgba(255,255,255,0.08);
  border: 1px solid rgba(255,255,255,0.1);
  color: rgba(255,255,255,0.6);
  width: 36px; height: 36px;
  border-radius: 8px;
  font-size: 1.1rem;
  cursor: pointer;
  z-index: 50;
  display: flex; align-items: center; justify-content: center;
  transition: all 0.2s;
}}
.fullscreen-btn:hover {{
  background: rgba(255,255,255,0.15);
  color: #fff;
}}

/* Responsive */
@media (max-width: 768px) {{
  .slide {{
    width: 95vw;
    height: 85vh;
    border-radius: 8px;
  }}
  .slide-inner {{
    padding: 24px 28px;
  }}
  .cover-title {{ font-size: 1.8rem; }}
  .slide-title {{ font-size: 1.3rem; }}
  .card-grid {{
    grid-template-columns: 1fr;
  }}
  .nav-arrow {{ display: none; }}
}}

/* Print */
@media print {{
  body {{ background: #fff; }}
  .slide {{
    position: relative !important;
    opacity: 1 !important;
    transform: none !important;
    width: 100% !important;
    height: auto !important;
    box-shadow: none !important;
    border: 1px solid #e0e0e0;
    page-break-after: always;
    margin-bottom: 20px;
    pointer-events: auto !important;
    display: block !important;
  }}
  .progress-bar, .nav-dots, .slide-counter,
  .nav-arrow, .fullscreen-btn {{ display: none !important; }}
  .slides-container {{
    display: block;
    width: 100%;
    height: auto;
  }}
  .slide-cover {{ color: #202124; background: #f0f0f0 !important; }}
  .slide-cover .slide-title {{ color: #202124; }}
  .slide-cover .slide-subtitle {{ color: #5f6368; }}
}}
</style>
</head>
<body>

<!-- Progress bar -->
<div class="progress-bar"><div class="progress-fill" id="progressFill"></div></div>

<!-- Slides -->
<div class="slides-container" id="slidesContainer">
{slides_joined}
</div>

<!-- Navigation -->
<div class="nav-dots" id="navDots">
{dots}
</div>
<div class="slide-counter" id="slideCounter">1 / {total}</div>
<button class="fullscreen-btn" id="fullscreenBtn" title="Fullscreen (F)">&#x26F6;</button>
<button class="nav-arrow prev" id="prevBtn" title="Previous">&#x2039;</button>
<button class="nav-arrow next" id="nextBtn" title="Next">&#x203A;</button>

<script>
(function() {{
  const slides = document.querySelectorAll('.slide');
  const dots = document.querySelectorAll('.dot');
  const counter = document.getElementById('slideCounter');
  const progress = document.getElementById('progressFill');
  const total = slides.length;
  let current = 0;

  function goTo(n) {{
    if (n < 0 || n >= total || n === current) return;
    slides[current].classList.remove('active');
    slides[current].classList.add(n > current ? 'prev' : '');
    slides[current].classList.remove(n < current ? 'prev' : '');
    current = n;
    slides[current].classList.remove('prev');
    slides[current].classList.add('active');
    dots.forEach((d, i) => d.classList.toggle('active', i === current));
    counter.textContent = (current + 1) + ' / ' + total;
    progress.style.width = ((current + 1) / total * 100) + '%';
    // Reset inactive slides
    slides.forEach((s, i) => {{
      if (i !== current) {{
        s.classList.remove('active');
        s.classList.toggle('prev', i < current);
      }}
    }});
  }}

  function next() {{ goTo(current + 1); }}
  function prev() {{ goTo(current - 1); }}

  // Keyboard
  document.addEventListener('keydown', function(e) {{
    if (e.key === 'ArrowRight' || e.key === ' ') {{ next(); e.preventDefault(); }}
    else if (e.key === 'ArrowLeft') {{ prev(); e.preventDefault(); }}
    else if (e.key === 'f' || e.key === 'F') {{
      if (!document.fullscreenElement) document.documentElement.requestFullscreen();
      else document.exitFullscreen();
    }}
    else if (e.key === 'Escape' && document.fullscreenElement) {{
      document.exitFullscreen();
    }}
  }});

  // Click navigation
  document.getElementById('nextBtn').addEventListener('click', next);
  document.getElementById('prevBtn').addEventListener('click', prev);

  // Dot navigation
  dots.forEach(d => {{
    d.addEventListener('click', () => goTo(parseInt(d.dataset.index)));
  }});

  // Fullscreen
  document.getElementById('fullscreenBtn').addEventListener('click', function() {{
    if (!document.fullscreenElement) document.documentElement.requestFullscreen();
    else document.exitFullscreen();
  }});

  // Touch/swipe
  let touchStartX = 0;
  let touchStartY = 0;
  document.addEventListener('touchstart', function(e) {{
    touchStartX = e.changedTouches[0].screenX;
    touchStartY = e.changedTouches[0].screenY;
  }});
  document.addEventListener('touchend', function(e) {{
    const dx = e.changedTouches[0].screenX - touchStartX;
    const dy = e.changedTouches[0].screenY - touchStartY;
    if (Math.abs(dx) > Math.abs(dy) && Math.abs(dx) > 50) {{
      if (dx < 0) next(); else prev();
    }}
  }});

  // Init
  slides[0].classList.add('active');
  progress.style.width = (1 / total * 100) + '%';
}})();
</script>
</body>
</html>'''

    return html


# ============================================================
# MAIN
# ============================================================

def main():
    # Determine input file
    input_file = INPUT_FILE
    if not input_file and len(sys.argv) > 1:
        input_file = sys.argv[1]

    if not input_file:
        print("Usage: python3 bp2html_converter.py <input.pptx|input.pdf>")
        sys.exit(1)

    if not os.path.exists(input_file):
        print(f"File not found: {input_file}")
        sys.exit(1)

    ext = os.path.splitext(input_file)[1].lower()

    print(f"Extracting content from: {input_file}")

    if ext == ".pptx":
        slides = extract_from_pptx(input_file)
    elif ext == ".pdf":
        slides = extract_from_pdf(input_file)
    else:
        print(f"Unsupported format: {ext}. Use .pptx or .pdf")
        sys.exit(1)

    print(f"Extracted {len(slides)} slides/pages")

    title = os.path.splitext(os.path.basename(input_file))[0]
    html = generate_html(slides, title)

    output_file = os.path.splitext(input_file)[0] + "_presentation.html"
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(html)

    print(f"HTML presentation generated: {output_file}")
    print(f"Open in browser to view. Use arrow keys to navigate, F for fullscreen.")


if __name__ == "__main__":
    main()
